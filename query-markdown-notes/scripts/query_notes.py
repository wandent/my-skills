#!/usr/bin/env python3
"""
Query Markdown Notes - Search and retrieve from knowledge base
Supports full-text search, project filtering, and summary generation
"""

import os
import re
import json
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple
import sys

# Multi-format support
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

KNOWLEDGE_BASE_PATH = Path(r"C:\Users\wandent\OneDrive - Microsoft\Documents\_CST\CST_FY26")
PREPROCESSED_DIR = Path.cwd() / ".knowledge_base_preprocessed"


class MarkdownNotesSearcher:
    def __init__(self):
        self.kb_path = KNOWLEDGE_BASE_PATH
        self.preprocessed_dir = PREPROCESSED_DIR
        self.preprocessed_dir.mkdir(exist_ok=True)
        self.file_cache = {}
        self.index = {}
        
    def preprocess_docx(self, file_path: Path) -> Optional[Path]:
        """Convert DOCX to markdown"""
        if DocxDocument is None:
            return None
        
        try:
            doc = DocxDocument(file_path)
            output_path = self.preprocessed_dir / f"{file_path.stem}.md"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"# {file_path.stem}\n\n")
                for para in doc.paragraphs:
                    if para.text.strip():
                        f.write(para.text + "\n")
                
                # Add tables
                for table in doc.tables:
                    f.write("\n| " + " | ".join(cell.text for cell in table.rows[0].cells) + " |\n")
                    f.write("| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |\n")
                    for row in table.rows[1:]:
                        f.write("| " + " | ".join(cell.text for cell in row.cells) + " |\n")
                    f.write("\n")
            
            return output_path
        except Exception as e:
            print(f"Error processing DOCX {file_path}: {e}")
            return None
    
    def preprocess_pptx(self, file_path: Path) -> Optional[Path]:
        """Convert PPTX to markdown"""
        if Presentation is None:
            return None
        
        try:
            prs = Presentation(file_path)
            output_path = self.preprocessed_dir / f"{file_path.stem}.md"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"# {file_path.stem}\n\n")
                for i, slide in enumerate(prs.slides, 1):
                    f.write(f"## Slide {i}\n\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            f.write(shape.text + "\n\n")
            
            return output_path
        except Exception as e:
            print(f"Error processing PPTX {file_path}: {e}")
            return None
    
    def preprocess_xlsx(self, file_path: Path) -> Optional[Path]:
        """Convert XLSX to markdown"""
        if openpyxl is None:
            return None
        
        try:
            wb = openpyxl.load_workbook(file_path)
            output_path = self.preprocessed_dir / f"{file_path.stem}.md"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"# {file_path.stem}\n\n")
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    f.write(f"## {sheet_name}\n\n")
                    
                    rows = list(sheet.iter_rows(values_only=True))
                    if rows:
                        headers = rows[0]
                        f.write("| " + " | ".join(str(h) if h else "" for h in headers) + " |\n")
                        f.write("| " + " | ".join(["---"] * len(headers)) + " |\n")
                        
                        for row in rows[1:]:
                            f.write("| " + " | ".join(str(v) if v else "" for v in row) + " |\n")
                    f.write("\n")
            
            return output_path
        except Exception as e:
            print(f"Error processing XLSX {file_path}: {e}")
            return None
    
    def preprocess_files(self):
        """Find and preprocess non-markdown files"""
        # Skip large system folders to improve performance
        skip_dirs = {'.venv', 'node_modules', '__pycache__', '.git', 'Resources'}
        
        for file_path in self.kb_path.rglob("*"):
            # Skip if in excluded directories
            if any(skip_dir in file_path.parts for skip_dir in skip_dirs):
                continue
                
            if file_path.is_file():
                try:
                    if file_path.suffix.lower() == ".docx":
                        self.preprocess_docx(file_path)
                    elif file_path.suffix.lower() == ".pptx":
                        self.preprocess_pptx(file_path)
                    elif file_path.suffix.lower() == ".xlsx":
                        self.preprocess_xlsx(file_path)
                except Exception as e:
                    # Silently skip files that can't be processed
                    pass
    
    def load_markdown_files(self, project_filter: Optional[str] = None) -> List[Path]:
        """Load markdown files from knowledge base"""
        pattern = "**/*.md"
        files = list(self.kb_path.glob(pattern))
        
        # Also include preprocessed files
        if self.preprocessed_dir.exists():
            files.extend(self.preprocessed_dir.glob("*.md"))
        
        if project_filter:
            # Filter by project path
            filtered = []
            for f in files:
                try:
                    rel_path = f.relative_to(self.kb_path)
                    if project_filter.lower() in str(rel_path).lower():
                        filtered.append(f)
                except ValueError:
                    # File is in preprocessed dir, skip project filter
                    pass
            return filtered
        
        return files
    
    def search_files(self, query: str, max_results: int = 10, 
                    context_lines: int = 3, project_filter: Optional[str] = None) -> List[Dict]:
        """Search files for query term"""
        # Pre-process any non-markdown files first
        self.preprocess_files()
        
        results = []
        files = self.load_markdown_files(project_filter)
        
        # Compile search patterns
        query_pattern = re.compile(re.escape(query), re.IGNORECASE)
        
        for file_path in files:
            try:
                # Skip .venv and other hidden dirs
                if ".venv" in str(file_path) or "node_modules" in str(file_path):
                    continue
                
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    lines = content.split('\n')
                
                # Find matches
                matches = []
                for i, line in enumerate(lines):
                    if query_pattern.search(line):
                        start = max(0, i - context_lines)
                        end = min(len(lines), i + context_lines + 1)
                        context = '\n'.join(lines[start:end])
                        
                        matches.append({
                            'line_num': i + 1,
                            'context': context,
                            'match_line': line
                        })
                
                if matches:
                    # Determine relative path
                    try:
                        rel_path = file_path.relative_to(self.kb_path)
                    except ValueError:
                        rel_path = file_path.relative_to(self.preprocessed_dir)
                    
                    # Extract project name
                    path_parts = str(rel_path).split(os.sep)
                    project = path_parts[1] if len(path_parts) > 1 else "Root"
                    
                    results.append({
                        'file': str(rel_path),
                        'project': project,
                        'match_count': len(matches),
                        'matches': matches[:3],  # Top 3 matches per file
                        'first_match_line': matches[0]['line_num']
                    })
            
            except Exception as e:
                continue
        
        # Sort by relevance (number of matches, then by line number)
        results.sort(key=lambda x: (-x['match_count'], x['first_match_line']))
        return results[:max_results]
    
    def list_projects(self, project_path: str = "Projects") -> Dict:
        """List available projects and structure"""
        path = self.kb_path / project_path
        
        if not path.exists():
            return {"error": f"Path not found: {project_path}"}
        
        structure = {}
        for item in sorted(path.iterdir()):
            if item.is_dir() and not item.name.startswith('.'):
                md_files = list(item.glob("**/*.md"))
                structure[item.name] = {
                    'type': 'folder',
                    'file_count': len(md_files),
                    'path': str(item.relative_to(self.kb_path))
                }
            elif item.is_file() and item.suffix == '.md':
                structure[item.name] = {
                    'type': 'file',
                    'path': str(item.relative_to(self.kb_path))
                }
        
        return structure
    
    def generate_summary(self, query: str, max_files: int = 5) -> str:
        """Generate summary from search results"""
        # Search for files
        results = self.search_files(query, max_results=max_files)
        
        summary_lines = [
            f"# Summary: {query}",
            "",
            "## Summary",
            f"Retrieved from knowledge base search for: **{query}**",
            ""
        ]
        
        # Add findings
        if results:
            summary_lines.append("## Key Findings")
            summary_lines.append("")
            
            for result in results:
                summary_lines.append(f"### {result['file']}")
                summary_lines.append(f"Project: {result['project']}")
                summary_lines.append("")
                
                # Add snippets
                for match in result['matches']:
                    summary_lines.append("**Relevant Section:**")
                    summary_lines.append("```")
                    summary_lines.append(match['context'][:300])  # Truncate
                    summary_lines.append("```")
                    summary_lines.append("")
            
            # Add sources
            summary_lines.append("## Sources")
            summary_lines.append("")
            for result in results:
                summary_lines.append(f"- {result['file']} ({result['match_count']} matches)")
            
            summary_lines.append("")
            summary_lines.append("*[Internal full context loaded for reasoning]*")
        else:
            summary_lines.append("No results found.")
        
        summary_content = '\n'.join(summary_lines)
        
        # Save to file
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = Path.cwd() / f"summary_{query.replace(' ', '_')}_{timestamp}.md"
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(summary_content)
        
        return str(output_file)


def main():
    """CLI interface for the search tool"""
    if len(sys.argv) < 2:
        print("Usage: query_notes.py <search_query> [--project <project>] [--max-results <n>] [--summary]")
        sys.exit(1)
    
    searcher = MarkdownNotesSearcher()
    query = sys.argv[1]
    
    # Parse arguments
    project_filter = None
    max_results = 10
    generate_summary = False
    
    i = 2
    while i < len(sys.argv):
        if sys.argv[i] == "--project" and i + 1 < len(sys.argv):
            project_filter = sys.argv[i + 1]
            i += 2
        elif sys.argv[i] == "--max-results" and i + 1 < len(sys.argv):
            max_results = int(sys.argv[i + 1])
            i += 2
        elif sys.argv[i] == "--summary":
            generate_summary = True
            i += 1
        else:
            i += 1
    
    # Perform search
    print(f"Searching for: {query}")
    if project_filter:
        print(f"In project: {project_filter}")
    print()
    
    results = searcher.search_files(query, max_results, project_filter=project_filter)
    
    if results:
        for result in results:
            print(f"[File: {result['file']}]")
            print(f"Project: {result['project']} ({result['match_count']} matches)")
            print()
            for match in result['matches']:
                print("Context:")
                print(match['context'])
                print()
            print("-" * 80)
            print()
    else:
        print("No results found.")
    
    # Generate summary if requested
    if generate_summary:
        summary_file = searcher.generate_summary(query)
        print(f"\nSummary saved to: {summary_file}")


if __name__ == "__main__":
    main()
